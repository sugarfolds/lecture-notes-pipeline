#!/usr/bin/env python3
from __future__ import annotations

import argparse
import re
import subprocess
from pathlib import Path

from pptx import Presentation


PPT_DIR = Path("ppt")
OUT_PATH = Path("slides_index.md")


def normalize_lines(raw_lines: list[str]) -> list[str]:
    lines: list[str] = []
    seen: set[str] = set()
    for raw in raw_lines:
        line = re.sub(r"\s+", " ", raw).strip()
        if not line or line in seen:
            continue
        seen.add(line)
        lines.append(line)
    return lines


def extract_pdf_lines(pdf: Path) -> list[str]:
    result = subprocess.run(
        ["pdftotext", str(pdf), "-"],
        check=True,
        capture_output=True,
        text=True,
    )
    return normalize_lines(result.stdout.replace("\f", "\n").splitlines())


def extract_pptx_lines(pptx_path: Path) -> list[str]:
    presentation = Presentation(pptx_path)
    raw_lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                raw_lines.extend(text.splitlines())
    return normalize_lines(raw_lines)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--ppt-dir", type=Path, default=PPT_DIR)
    parser.add_argument("--output", type=Path, default=OUT_PATH)
    args = parser.parse_args()

    course_name = args.ppt_dir.parent.name or args.ppt_dir.name
    rows = [f"# {course_name}课件索引", ""]
    files = sorted(args.ppt_dir.glob("*.pptx")) + sorted(args.ppt_dir.glob("*.pdf"))
    for path in files:
        if path.suffix.lower() == ".pptx":
            lines = extract_pptx_lines(path)
        else:
            lines = extract_pdf_lines(path)
        preview = lines[:10]
        rows.append(f"## {path.name}")
        rows.append("")
        for line in preview:
            rows.append(f"- {line}")
        rows.append("")
    args.output.write_text("\n".join(rows), encoding="utf-8")


if __name__ == "__main__":
    main()
